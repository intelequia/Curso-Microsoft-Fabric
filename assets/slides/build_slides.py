"""Genera las presentaciones PowerPoint del curso (Jornada 1 y Jornada 2).

Uso:
    python3 build_slides.py

Genera:
    J1-microsoft-fabric.pptx
    J2-purview-fabric-iq.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT_DIR = Path(__file__).resolve().parent

# Paleta tipo Fabric / Purview
PRIMARY = RGBColor(0x11, 0x2A, 0x55)      # azul oscuro
ACCENT = RGBColor(0x00, 0x78, 0xD4)       # azul Fabric
ACCENT2 = RGBColor(0x74, 0x2E, 0xCD)      # morado Purview / IQ
LIGHT = RGBColor(0xF3, 0xF6, 0xFB)
DARK_TEXT = RGBColor(0x20, 0x23, 0x2A)
MUTED = RGBColor(0x60, 0x66, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=18, color=DARK_TEXT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + b
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def title_slide(prs, eyebrow, title, subtitle, accent=ACCENT):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, accent)
    add_rect(slide, Inches(11.0), Inches(6.4), Inches(2.0), Inches(0.12), accent)
    add_rect(slide, Inches(11.6), Inches(6.65), Inches(1.4), Inches(0.06), WHITE)
    add_text(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.6),
             eyebrow, size=20, bold=True, color=accent)
    add_text(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2),
             title, size=54, bold=True, color=WHITE)
    add_text(slide, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.0),
             subtitle, size=22, color=LIGHT)
    add_text(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
             "Curso Microsoft Fabric · Caso ficticio: Aurora Energía",
             size=14, color=LIGHT)
    return slide


def section_slide(prs, eyebrow, title, accent=ACCENT):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT)
    add_rect(slide, 0, Inches(2.6), SLIDE_W, Inches(2.3), PRIMARY)
    add_rect(slide, 0, Inches(2.6), Inches(0.35), Inches(2.3), accent)
    add_text(slide, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.6),
             eyebrow, size=18, bold=True, color=accent)
    add_text(slide, Inches(0.9), Inches(3.35), Inches(11.5), Inches(1.4),
             title, size=44, bold=True, color=WHITE)
    return slide


def content_slide(prs, eyebrow, title, bullets, *, footer=None, accent=ACCENT,
                  body_size=18):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), PRIMARY)
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.08), accent)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.4),
             eyebrow, size=14, bold=True, color=accent)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.5), Inches(0.8),
             title, size=28, bold=True, color=WHITE)
    add_bullets(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.2),
                bullets, size=body_size)
    if footer:
        add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), LIGHT)
        add_text(slide, Inches(0.5), Inches(7.10), Inches(12.5), Inches(0.4),
                 footer, size=11, color=MUTED)
    return slide


def two_col_slide(prs, eyebrow, title, left_title, left_bullets,
                  right_title, right_bullets, *, accent=ACCENT, footer=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), PRIMARY)
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.08), accent)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.4),
             eyebrow, size=14, bold=True, color=accent)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.5), Inches(0.8),
             title, size=28, bold=True, color=WHITE)
    col_w = Inches(6.0)
    col_h = Inches(5.4)
    left_x = Inches(0.5)
    right_x = Inches(6.85)
    add_rect(slide, left_x, Inches(1.55), col_w, Inches(0.55), accent)
    add_text(slide, left_x + Inches(0.2), Inches(1.62), col_w, Inches(0.45),
             left_title, size=18, bold=True, color=WHITE)
    add_rect(slide, left_x, Inches(2.10), col_w, col_h - Inches(0.55), WHITE)
    add_bullets(slide, left_x + Inches(0.2), Inches(2.20), col_w - Inches(0.4),
                col_h - Inches(0.7), left_bullets, size=15)
    add_rect(slide, right_x, Inches(1.55), col_w, Inches(0.55), ACCENT2)
    add_text(slide, right_x + Inches(0.2), Inches(1.62), col_w, Inches(0.45),
             right_title, size=18, bold=True, color=WHITE)
    add_rect(slide, right_x, Inches(2.10), col_w, col_h - Inches(0.55), WHITE)
    add_bullets(slide, right_x + Inches(0.2), Inches(2.20), col_w - Inches(0.4),
                col_h - Inches(0.7), right_bullets, size=15)
    if footer:
        add_text(slide, Inches(0.5), Inches(7.10), Inches(12.5), Inches(0.4),
                 footer, size=11, color=MUTED)
    return slide


def quote_slide(prs, eyebrow, quote, attrib=None, accent=ACCENT):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT)
    add_rect(slide, Inches(0.7), Inches(1.4), Inches(0.18), Inches(4.7), accent)
    add_text(slide, Inches(1.1), Inches(1.4), Inches(11.4), Inches(0.6),
             eyebrow, size=16, bold=True, color=accent)
    add_text(slide, Inches(1.1), Inches(2.0), Inches(11.4), Inches(4.0),
             "“" + quote + "”", size=30, bold=True, color=PRIMARY)
    if attrib:
        add_text(slide, Inches(1.1), Inches(6.0), Inches(11.4), Inches(0.5),
                 "— " + attrib, size=16, color=MUTED)
    return slide


def agenda_slide(prs, eyebrow, title, rows, accent=ACCENT):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), PRIMARY)
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.08), accent)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.4),
             eyebrow, size=14, bold=True, color=accent)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.5), Inches(0.8),
             title, size=28, bold=True, color=WHITE)
    y = Inches(1.6)
    row_h = Inches(0.46)
    x_code = Inches(0.7)
    x_dur = Inches(2.2)
    x_topic = Inches(3.7)
    add_rect(slide, x_code - Inches(0.1), y, Inches(12.0), row_h, accent)
    add_text(slide, x_code, y + Inches(0.05), Inches(1.4), row_h,
             "Bloque", size=14, bold=True, color=WHITE)
    add_text(slide, x_dur, y + Inches(0.05), Inches(1.4), row_h,
             "Duración", size=14, bold=True, color=WHITE)
    add_text(slide, x_topic, y + Inches(0.05), Inches(8.6), row_h,
             "Tema", size=14, bold=True, color=WHITE)
    for i, (code, dur, topic) in enumerate(rows):
        ry = y + row_h * (i + 1)
        bg = WHITE if i % 2 == 0 else LIGHT
        add_rect(slide, x_code - Inches(0.1), ry, Inches(12.0), row_h, bg)
        add_text(slide, x_code, ry + Inches(0.07), Inches(1.4), row_h,
                 code, size=13, bold=True, color=PRIMARY)
        add_text(slide, x_dur, ry + Inches(0.07), Inches(1.4), row_h,
                 dur, size=13, color=MUTED)
        add_text(slide, x_topic, ry + Inches(0.07), Inches(8.6), row_h,
                 topic, size=13, color=DARK_TEXT)
    return slide


def closing_slide(prs, title, bullets, accent=ACCENT):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, accent)
    add_text(slide, Inches(0.9), Inches(0.9), Inches(11.5), Inches(1.0),
             title, size=42, bold=True, color=WHITE)
    add_bullets(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(4.5),
                bullets, size=22, color=LIGHT)
    add_text(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
             "¡Gracias!  ·  Curso Microsoft Fabric + Purview + Fabric IQ",
             size=14, color=LIGHT)
    return slide


# =========================================================================
# JORNADA 1 — Microsoft Fabric
# =========================================================================

def build_jornada_1():
    prs = new_presentation()
    A = ACCENT

    title_slide(
        prs,
        "JORNADA 1 · 4 horas",
        "Microsoft Fabric",
        "De cero a un flujo end-to-end: ingesta, transformación, modelo y reporte",
        accent=A,
    )

    agenda_slide(
        prs, "Agenda · Jornada 1", "Cómo vamos a invertir las próximas 4 horas",
        [
            ("J1-M0", "15 min", "Bienvenida, agenda y caso Aurora Energía"),
            ("J1-M1", "30 min", "Fundamentos: arquitectura, capacidades, OneLake"),
            ("J1-M2", "35 min", "Lakehouse, Warehouse y SQL Endpoint"),
            ("J1-M3", "40 min", "Ingesta: Dataflow Gen2 y Data Pipelines"),
            ("—",     "15 min", "Descanso"),
            ("J1-M4", "35 min", "Notebooks, Spark y procesamiento"),
            ("J1-M5", "30 min", "Warehouse en profundidad y Direct Lake"),
            ("J1-M6", "25 min", "Real-Time Intelligence: Eventhouse y KQL"),
            ("J1-M7", "25 min", "Power BI sobre Fabric: modelo semántico"),
            ("J1-M8", "10 min", "Cierre, dudas y deberes para casa"),
        ], accent=A,
    )

    content_slide(
        prs, "J1-M0 · Bienvenida", "Caso conductor: Aurora Energía",
        [
            "Operador ficticio: red de estaciones de servicio, comercializadora eléctrica y división de logística.",
            "Todos los datos, esquemas y nombres son sintéticos — sin referencias a clientes reales.",
            "Construiremos incrementalmente el mismo entorno (workspace aurora-curso-fabric).",
            "Metodología: explicación + demo en vivo + 2 ejercicios cortos en aula + pack de prácticas para casa.",
            "Lo construido en clase es la base del trabajo de Jornada 2 (Purview + Fabric IQ).",
        ], accent=A,
    )

    section_slide(prs, "J1-M1 · 30 min", "Fundamentos de Microsoft Fabric", accent=A)
    content_slide(
        prs, "J1-M1 · Fundamentos", "¿Qué es Microsoft Fabric?",
        [
            "Plataforma SaaS de analítica unificada (GA noviembre 2023).",
            "Reúne ingesta, data engineering, warehouse, real-time, data science, Power BI y, desde 2024-26, Fabric Databases y Fabric IQ.",
            "Tres pilares: OneLake · Capacidades F · Experiencias por rol.",
            "SaaS-first: no se aprovisiona infraestructura, se crea workspace y se empieza.",
            "Una sola capacidad de cómputo unificada para Spark, T-SQL, KQL y Power BI.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M1 · Fundamentos", "OneLake — el OneDrive de los datos",
        [
            "Un único Data Lake lógico por tenant, multi-cloud, sobre ADLS Gen2.",
            "Jerarquía: Tenant → Workspace → Item (Lakehouse / Warehouse / Eventhouse) → Files & Tables.",
            "Todo en formato abierto Delta-Parquet → cualquier motor que lea Delta puede leer OneLake.",
            "Shortcuts: punteros virtuales a datos internos o externos (ADLS, S3, GCS) sin duplicar.",
            "Mirroring: réplica casi en tiempo real de Cosmos DB, Azure SQL DB, Snowflake, Fabric SQL DB…",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M1 · Fundamentos", "Workloads / Experiencias",
        [
            "Data Factory → ingesta (Pipelines, Dataflow Gen2, Copy Job).",
            "Data Engineering → Lakehouse, Notebooks, Spark Jobs, Environments.",
            "Data Warehouse → Warehouse T-SQL nativo.",
            "Data Science → notebooks, MLflow, modelos.",
            "Real-Time Intelligence → Eventstream, Eventhouse (KQL DB), Activator.",
            "Power BI → semantic models, reports, dashboards, apps.",
            "Industry Solutions y Fabric Databases (SQL DB).",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M1 · Fundamentos", "Capacidades, licencias y trial",
        [
            "F-SKU (F2 → F2048): unidad de cómputo, facturada por hora, pausable. F64 ≈ P1 de Power BI Premium.",
            "Power BI Free / Pro / PPU: necesarios para consumir reportes si la capacidad es < F64.",
            "Regla de oro: para distribuir reportes a usuarios Free necesitas mínimo F64.",
            "Trial: 60 días, equivalente a FT1 (~F64), sin tarjeta, una por usuario.",
            "Roles del workspace: Admin · Member · Contributor · Viewer (heredados de Power BI).",
        ], accent=A, footer="Demo: portal Fabric · OneLake Catalog · shortcuts · capacidad asignada al workspace.",
    )
    quote_slide(
        prs, "Mensaje clave",
        "OneLake es la pieza realmente disruptiva: un único almacenamiento, una única gobernanza, una sola capacidad de cómputo.",
        accent=A,
    )

    section_slide(prs, "J1-M2 · 35 min", "Lakehouse, Warehouse y SQL Endpoint", accent=A)
    content_slide(
        prs, "J1-M2 · Lakehouse", "Lakehouse en Fabric",
        [
            "Estructura física: zona Files (cualquier formato) + zona Tables (Delta-Parquet).",
            "Ingesta multi-modo: portal, OneLake File Explorer, Spark, Dataflow, Pipeline, shortcut.",
            "Cada Lakehouse trae GRATIS un SQL Analytics Endpoint (solo lectura) y un Default Semantic Model.",
            "Encaje ideal: data lake moderno, ingesta cruda + curado, ML, Spark, datos no estructurados.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M2 · Warehouse", "Warehouse en Fabric",
        [
            "T-SQL nativo, lectura/escritura: INSERT, UPDATE, DELETE, MERGE.",
            "Almacena en Delta-Parquet sobre OneLake (mismo formato que el Lakehouse).",
            "Soporta vistas, procs, funciones, transacciones cross-table, RLS, CLS, OLS, masking.",
            "Encaje ideal: equipo SQL puro, modelo en estrella servido a Power BI, lógica con DML/transacciones.",
        ], accent=A,
    )
    two_col_slide(
        prs, "J1-M2 · Comparativa", "Lakehouse vs Warehouse",
        "Lakehouse (+ SQL endpoint)",
        [
            "Lenguaje principal: PySpark / Spark SQL",
            "Escritura SQL: ❌ (endpoint solo lectura)",
            "Datos no estructurados: ✅",
            "DML T-SQL: ❌",
            "Transacciones multi-tabla: limitadas",
            "RLS / CLS / OLS: parcial",
            "Power BI Direct Lake: ✅",
        ],
        "Warehouse",
        [
            "Lenguaje principal: T-SQL",
            "Escritura SQL: ✅",
            "Datos no estructurados: ❌",
            "DML T-SQL: ✅",
            "Transacciones multi-tabla: ✅ ACID",
            "RLS / CLS / OLS: completo",
            "Power BI Direct Lake: ✅",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M2 · Decisión", "¿Cuándo elijo cada uno?",
        [
            "Lakehouse: ingesta heterogénea, ciencia de datos, no necesito DML SQL.",
            "Warehouse: equipo SQL puro, DML/MERGE, transacciones y gobierno fino.",
            "Patrón habitual: Lakehouse para bronze/silver + Warehouse para gold servido a Power BI.",
            "Caso Aurora: lh_aurora (landing) + wh_aurora (modelo dimensional).",
            "Modelo Direct Lake apuntará al Warehouse en el bloque M7.",
        ], accent=A, footer="Demo + ejercicio (10 min): crear lh_aurora, wh_aurora, cargar clientes.csv, dim_cliente vacía.",
    )

    section_slide(prs, "J1-M3 · 40 min", "Ingesta: Dataflow Gen2 y Data Pipelines", accent=A)
    content_slide(
        prs, "J1-M3 · Ingesta", "Opciones de ingesta en Fabric",
        [
            "Dataflow Gen2: ingesta + transformación con Power Query, low-code, batch (M).",
            "Data Pipeline: orquestación, control de flujo, copia masiva, programación (UI).",
            "Copy Job (GA reciente): copia incremental gestionada (CDC + watermark).",
            "Eventstream: streaming continuo desde IoT Hub, Event Hubs, Kafka, sample data.",
            "Mirroring: réplica casi en tiempo real desde Azure SQL, Cosmos, Snowflake, Fabric SQL DB.",
            "Notebook Spark: ingesta programática (spark.read).",
        ], accent=A,
    )
    two_col_slide(
        prs, "J1-M3 · Comparativa", "Dataflow Gen2 vs Data Pipeline",
        "Dataflow Gen2",
        [
            "Power Query Online (sin código).",
            "Multi-destino: Lakehouse, Warehouse, KQL DB, Azure SQL.",
            "Plantillas .pqt importables/exportables.",
            "Fast Copy en orígenes/conectores grandes.",
            "Refresco programable.",
        ],
        "Data Pipeline",
        [
            "Heredero de ADF v2 dentro de Fabric.",
            "Actividades: Copy, Dataflow, Notebook, Stored Proc, Lookup, ForEach…",
            "Eventos: éxito, fallo, omisión, completado.",
            "Triggers: cron, file arrived, manual.",
            "Monitoring Hub cross-workspace.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M3 · Patrón", "Pipeline orquesta el Dataflow",
        [
            "pl_aurora_ingesta → df_clientes → df_ventas → nb_silver_clean → sp_load_dim_cliente.",
            "Rama de éxito: notificar por Office365 / Teams.",
            "Rama de fallo: alerta al canal de soporte.",
            "Mirroring para orígenes operacionales (sin construir ETL).",
            "Eventstream para streaming continuo → Eventhouse / Lakehouse.",
        ], accent=A, footer="Demo (15 min) + ejercicio (5 min): df_clientes con destino lh_aurora.clientes.",
    )
    quote_slide(prs, "Mensaje clave",
                "Dataflow para transformar, Pipeline para orquestar. Si el origen es operacional, no construyas ETL: usa Mirroring.",
                accent=A)

    section_slide(prs, "Descanso · 15 min", "Estiramos las piernas y volvemos", accent=A)

    section_slide(prs, "J1-M4 · 35 min", "Notebooks, Spark y procesamiento de datos", accent=A)
    content_slide(
        prs, "J1-M4 · Notebooks", "Notebooks en Fabric",
        [
            "Editor estilo Jupyter con celdas de código y markdown.",
            "Lenguajes: PySpark, Spark SQL (%%sql), Scala, SparkR / sparklyr.",
            "Conexión nativa a uno o varios Lakehouses desde el explorador lateral.",
            "Spark gestionado: Starter Pool listo en segundos, sin provisionar cluster.",
            "Custom pools, Environments versionables (paquetes + Spark conf), sesiones high-concurrency.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M4 · APIs", "PySpark, Spark SQL y NotebookUtils",
        [
            "spark.read.format(\"csv\") · df.write.format(\"delta\").saveAsTable(...) · df.groupBy().agg().",
            "Celda %%sql para alumnos que vienen de SQL puro.",
            "NotebookUtils (mssparkutils / notebookutils): OneLake, secretos, fs, jobs encadenados.",
            "Pandas API on Spark y fabric-data-functions para lectura pandas-like.",
            "display(df) en lugar de df.show() para visuales interactivos.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M4 · Caso Aurora", "Notebook nb_aurora_lab: bronze → silver",
        [
            "Lee clientes, productos, estaciones y ventas_raw del Lakehouse.",
            "Limpia ventas_raw: quita importes ≤ 0, normaliza fechas, deriva año, mes, día_semana.",
            "Join con productos para añadir categoría y unidad de medida.",
            "Escribe ventas_silver como tabla Delta.",
            "Buenas prácticas: Git integration, Spark Job Definition para producción, particionado Delta.",
        ], accent=A, footer="Demo (15 min) + ejercicio para casa: ej04-notebook-bronze-silver.md.",
    )

    section_slide(prs, "J1-M5 · 30 min", "Warehouse en profundidad y Direct Lake", accent=A)
    content_slide(
        prs, "J1-M5 · Warehouse", "Modelado del gold (estrella) en Aurora",
        [
            "Dimensiones: dim_cliente, dim_producto, dim_estacion, dim_tiempo.",
            "Hechos: fact_ventas (venta_id, fecha, cliente_id, producto_id, estacion_id, cantidad, importe).",
            "Carga vía CTAS, SELECT INTO, MERGE pattern y stored procedures.",
            "Cross-database query desde Warehouse al Lakehouse (lh_aurora.dbo.clientes).",
            "Vistas y funciones para servir KPIs estables a Power BI.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M5 · Seguridad SQL", "RLS, CLS, OLS y masking",
        [
            "Row-Level Security: política basada en USER_NAME() o claims de Entra.",
            "Column-Level Security: GRANT SELECT solo sobre columnas concretas.",
            "Object-Level Security: ocultar tablas o columnas por rol.",
            "Dynamic Data Masking para columnas sensibles (DNI, email).",
            "En Jornada 2 lo cruzaremos con Purview labels.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M5 · Direct Lake", "El santo grial del modelo semántico",
        [
            "Power BI lee directamente las tablas Delta del Lakehouse / Warehouse.",
            "Velocidad de Import + frescura de DirectQuery, sin coste de refresh.",
            "Fallback automático a DirectQuery si la consulta excede límites.",
            "Requiere capacidad F-SKU, columnas tipadas y modelo creado desde el item.",
            "Direct Lake on OneLake (DLOL): modelo PBIP independiente del Lakehouse origen.",
        ], accent=A, footer="Demo (12 min): DDL + sp_load_dim_cliente + cross-DB + nuevo modelo Direct Lake.",
    )

    section_slide(prs, "J1-M6 · 25 min", "Real-Time Intelligence: Eventhouse y KQL", accent=A)
    content_slide(
        prs, "J1-M6 · Real-Time", "Mapa de la Real-Time Intelligence",
        [
            "Eventstream: ingesta de streaming, no-code, transformaciones simples.",
            "Eventhouse / KQL Database: motor analítico para series temporales y logs.",
            "Activator: motor de reglas que dispara correo, Teams, Power Automate o Pipelines.",
            "Real-Time Dashboard: dashboard tipo Grafana/Kusto Explorer integrado.",
            "Caso Aurora: 250 estaciones × 6 surtidores, eventos cada 30 s.",
        ], accent=A,
    )
    two_col_slide(
        prs, "J1-M6 · Decisión", "Eventhouse vs Lakehouse / Warehouse",
        "Eventhouse / KQL",
        [
            "Telemetría IoT, logs, clickstream, security events.",
            "Query en milisegundos sobre miles de millones de eventos.",
            "Funciones nativas de series temporales (make-series, anomalías).",
        ],
        "Lakehouse / Warehouse",
        [
            "Histórico analítico de negocio (ventas, finanzas).",
            "Unión a modelo dimensional.",
            "Servir Power BI con Direct Lake.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M6 · KQL en 5 minutos", "Sintaxis pipe-based",
        [
            "TelemetriaSurtidor | where Timestamp > ago(24h) | summarize count() by EstacionId | top 10 by count_ desc",
            "summarize avg(Caudal) by bin(Timestamp, 5m), EstacionId — agregaciones por ventana.",
            "make-series + series_decompose_anomalies para detección de anomalías.",
            "Activator: ej. \"Si Temperatura > 65 °C durante > 3 min, avisar al jefe de estación\".",
        ], accent=A, footer="Demo (10 min): eh_aurora_telemetria + Eventstream + KQL + Real-Time Dashboard + Activator.",
    )

    section_slide(prs, "J1-M7 · 25 min", "Power BI sobre Fabric: modelo semántico", accent=A)
    content_slide(
        prs, "J1-M7 · Modelo semántico", "Default vs Custom",
        [
            "Default semantic model del Lakehouse: tablas autodetectadas, sin relaciones definidas.",
            "Custom semantic model: relaciones, jerarquías, medidas DAX, perspectivas, RLS.",
            "Recomendación de producción: SIEMPRE custom semantic model.",
            "Se crea desde el item Lakehouse / Warehouse o desde Power BI Desktop.",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M7 · DAX", "Medidas típicas para Aurora",
        [
            "Importe Total = SUM(fact_ventas[importe]).",
            "Litros Vendidos = SUM(fact_ventas[cantidad]).",
            "Importe Año Anterior = CALCULATE([Importe Total], SAMEPERIODLASTYEAR(dim_tiempo[fecha])).",
            "Variación % = DIVIDE([Importe Total]-[Importe Año Anterior], [Importe Año Anterior]).",
            "Top 5 Estaciones = CALCULATE([Importe Total], TOPN(5, dim_estacion, [Importe Total])).",
        ], accent=A,
    )
    content_slide(
        prs, "J1-M7 · Producto", "Reportes, Apps, PBIP y Git",
        [
            "Reporte web directo en Fabric o publicación desde Power BI Desktop.",
            "Apps para empaquetar y distribuir (lectores Free necesitan F64+).",
            "Subscriptions, alerts, comments, Workspace OneLake hub.",
            "Power BI Desktop guarda en formato PBIP (carpeta + JSON/TMDL/PBIR), versionable en Git.",
            "Source control integrado en workspaces Fabric (Azure DevOps / GitHub).",
        ], accent=A, footer="Demo (12 min): sm_aurora_ventas + reporte + App + Source control.",
    )

    content_slide(
        prs, "J1-M8 · Cierre", "Recapitulación: el flujo end-to-end",
        [
            "CSV/Excel → Dataflow Gen2 → Lakehouse (bronze)",
            "→ Notebook Spark → Lakehouse (silver)",
            "→ Stored Proc Warehouse → Warehouse (gold)",
            "→ Direct Lake → Semantic Model → Reporte Power BI",
            "En paralelo: Event Hub → Eventstream → Eventhouse → KQL Dashboard / Activator",
        ], accent=A,
    )
    closing_slide(
        prs, "Deberes para casa",
        [
            "Completar todos los ejercicios de ejercicios/jornada-1/.",
            "Workspace aurora-curso-fabric con Lakehouse, Warehouse, Notebook, Pipeline y Reporte publicado.",
            "Lectura ligera del índice de ejercicios/jornada-2/ para llegar contextualizado.",
            "Canal de soporte abierto entre jornadas.",
            "Nos vemos en la Jornada 2: Microsoft Purview + Fabric IQ.",
        ], accent=A,
    )

    out = OUT_DIR / "J1-microsoft-fabric.pptx"
    prs.save(out)
    return out


# =========================================================================
# JORNADA 2 — Microsoft Purview + Fabric IQ
# =========================================================================

def build_jornada_2():
    prs = new_presentation()
    A = ACCENT2

    title_slide(
        prs,
        "JORNADA 2 · 4 horas",
        "Microsoft Purview + Fabric IQ",
        "Gobierno extremo a extremo y la nueva capa de inteligencia sobre el dato",
        accent=A,
    )

    agenda_slide(
        prs, "Agenda · Jornada 2", "Cómo vamos a invertir las próximas 4 horas",
        [
            ("J2-M0", "15 min", "Repaso Jornada 1, agenda y dudas"),
            ("J2-M1", "30 min", "Por qué gobierno: panorama Purview + Fabric"),
            ("J2-M2", "40 min", "Purview Data Map y Unified Catalog"),
            ("J2-M3", "30 min", "Information Protection, sensitivity labels y DLP"),
            ("—",     "15 min", "Descanso"),
            ("J2-M4", "25 min", "Integración Purview ↔ Fabric"),
            ("J2-M5", "35 min", "Fabric IQ: qué es, arquitectura y posicionamiento"),
            ("J2-M6", "30 min", "Fabric IQ en acción: Data Agents y Q&A"),
            ("J2-M7", "25 min", "Casos de uso, mejores prácticas y roadmap"),
            ("J2-M8", "15 min", "Cierre, próximos pasos y recursos"),
        ], accent=A,
    )

    content_slide(
        prs, "J2-M0 · Bienvenida", "De qué hablaremos hoy",
        [
            "Ayer construimos el dato en Fabric. Hoy lo gobernamos y lo hacemos accionable por IA.",
            "Bloque 1 (Purview): Data Map, Unified Catalog, sensitivity labels, DLP.",
            "Bloque 2 (Integración): Purview ↔ Fabric extremo a extremo.",
            "Bloque 3 (Fabric IQ): capa semántica, Data Agents y Q&A en lenguaje natural.",
            "Cierre: hoja de ruta de adopción a 90 días para Aurora Energía.",
        ], accent=A,
    )

    section_slide(prs, "J2-M1 · 30 min", "Por qué gobierno: panorama Microsoft Purview", accent=A)
    content_slide(
        prs, "J2-M1 · Purview", "Mapa de Microsoft Purview (2026)",
        [
            "Data Governance: Data Map, Unified Catalog, Data Quality, Lineage, Data Products, Access Policies.",
            "Information Protection: sensitivity labels, auto-labeling, encriptación, RMS.",
            "Risk & Compliance: DLP, Insider Risk, Communication Compliance, eDiscovery, Audit, Compliance Manager.",
            "Cambio de marca: \"Microsoft 365 Compliance Center\" → Purview Compliance Portal.",
            "Y \"Azure Purview\" → Purview Data Governance. Mismo portal, dos áreas.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M1 · Purview", "¿Por qué gobernar Fabric?",
        [
            "Fabric centraliza datos antes dispersos → el riesgo de exposición sube.",
            "OneLake permite que cualquier workspace lea/escriba → disciplina o caos.",
            "Compliance (GDPR, ENS, ISO 27001, NIS2) exige clasificar, rastrear y proteger.",
            "El negocio necesita encontrar el dato: un Catálogo es la única forma escalable.",
            "Sin gobierno, Fabric escala como cualquier data lake: hacia el caos.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M1 · Purview", "Necesidad → módulo Purview",
        [
            "¿Qué datos tengo y dónde? → Data Map + Unified Catalog.",
            "¿Quién tiene acceso? → Data Access Policies + integración Fabric.",
            "¿Es de calidad este dataset? → Data Quality.",
            "¿De dónde viene esta tabla? → Data Lineage.",
            "¿Qué columnas son PII / financieras? → Sensitivity labels + Auto-classification.",
            "¿Cómo cumplo GDPR / borrado? → Records Management + DSR + eDiscovery.",
            "¿Quién descarga datasets sensibles? → Insider Risk + Audit.",
        ], accent=A, footer="Demo (8 min): purview.microsoft.com · Unified Catalog · lineage de lh_aurora.",
    )
    quote_slide(prs, "Mensaje clave",
                "Purview es un paraguas. Para Fabric, no necesitas hacer scan: la integración nativa publica todo automáticamente.",
                accent=A)

    section_slide(prs, "J2-M2 · 40 min", "Purview Data Map y Unified Catalog", accent=A)
    two_col_slide(
        prs, "J2-M2 · Capas", "Data Map vs Unified Catalog",
        "Data Map (capa técnica)",
        [
            "Inventario de assets: tablas, ficheros, modelos, dashboards.",
            "Conexiones: Fabric (auto), Azure SQL, ADLS, Databricks, Snowflake, Power BI.",
            "On-prem vía Self-Hosted Integration Runtime.",
            "Scans programables con clasificaciones automáticas.",
        ],
        "Unified Catalog (capa de negocio)",
        [
            "Vista para analistas y owners de dato.",
            "Business Domain → Data Product → Data Asset.",
            "Glossary terms, Critical Data Elements, OKRs.",
            "Endorsements: Promoted ✅, Certified ⭐.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M2 · Calidad y linaje", "Data Quality, Lineage y Access Policies",
        [
            "Data Quality: reglas de completeness, uniqueness, validity, accuracy, freshness — score histórico.",
            "Data Lineage automático en Fabric: Lakehouse → Notebook → Warehouse → Semantic Model → Report.",
            "Lineage manual / programático con APIs y OpenLineage para fuentes externas.",
            "Data Access Policies: conceder acceso desde Purview sin entrar a la herramienta origen.",
            "Federa ADLS Gen2, Azure SQL DB y Fabric.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M2 · Caso Aurora", "Data Product \"Ventas Aurora — Gold\"",
        [
            "Dominio: Comercial.",
            "Items asociados: wh_aurora + sm_aurora_ventas.",
            "Glossary terms: Importe de Venta · Estación de Servicio · Cliente Comercial.",
            "Critical Data Elements: cliente_id, importe.",
            "Owner: Marisa Ledesma (ficticia) · Endorsement: Certified.",
        ], accent=A, footer="Demo (15 min) + ejercicio (8 min): glossary term \"Importe de Venta\" asociado a fact_ventas[importe].",
    )

    section_slide(prs, "J2-M3 · 30 min", "Information Protection, etiquetas y DLP", accent=A)
    content_slide(
        prs, "J2-M3 · Sensitivity labels", "Ciclo de vida de las etiquetas",
        [
            "Definidas a nivel tenant desde Purview Compliance Portal.",
            "Pueden marcar visualmente, encriptar (MIP / Azure RMS) y restringir acciones.",
            "Aplican a Office, Outlook, Teams, SharePoint, OneDrive, items de Fabric, OneLake y Power BI.",
            "Aplicación: manual, auto-label por contenido, auto-label por contenedor.",
            "En Fabric: si etiqueto un Lakehouse, modelos y reportes derivados HEREDAN la etiqueta.",
            "Al exportar a Excel/PDF, se exporta también la etiqueta y la encriptación.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M3 · Clasificación", "Sensitive Information Types y EDM",
        [
            "Catálogo enorme out-of-the-box (DNI España, NIF, IBAN, tarjetas, pasaportes…).",
            "Personalizables: regex + keywords + función + nivel de confianza.",
            "EDM (Exact Data Match): comparar contra una tabla maestra propia.",
            "Trainable classifiers: ML para categorizar contratos, CVs, código fuente, etc.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M3 · Etiquetas Aurora", "Las 5 etiquetas de Aurora Energía",
        [
            "Público — sin marca ni encriptación (marketing).",
            "Interno — marca \"Aurora — Interno\" (documentación operativa).",
            "Confidencial — Comercial — encriptado para grupo Ventas.",
            "Restringido — PII — encriptado solo para Owner del Data Product.",
            "Restringido — Financiero — encriptado para grupo Finanzas.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M3 · DLP", "Data Loss Prevention en y fuera de Fabric",
        [
            "Reglas que detectan, alertan o bloquean contenidos sensibles.",
            "Alcance: Endpoint DLP, Exchange, Teams, SharePoint, OneDrive, Fabric/Power BI, terceros (Defender for Cloud Apps).",
            "Acciones: aviso, justificación obligatoria, bloqueo, override con auditoría.",
            "Ejemplo Fabric: \"Si un semantic model contiene > X registros con DNI, bloquea export a CSV\".",
            "Activity Explorer y Audit: visor cronológico para investigación.",
        ], accent=A, footer="Demo (10 min): aplicar Aurora — Confidencial a wh_aurora · regla DLP que bloquea export.",
    )
    quote_slide(prs, "Mensaje clave",
                "La etiqueta viaja con el dato — incluso cuando sale de Fabric. DLP no se inventa el día del incidente: se diseña antes.",
                accent=A)

    section_slide(prs, "Descanso · 15 min", "Estiramos las piernas y volvemos", accent=A)

    section_slide(prs, "J2-M4 · 25 min", "Integración Purview ↔ Fabric", accent=A)
    content_slide(
        prs, "J2-M4 · Vistas", "¿Dónde se ve qué?",
        [
            "Workspace Fabric → items técnicos del proyecto (equipo técnico).",
            "OneLake Catalog (en Fabric) → items del tenant filtrables, búsqueda, endorsements (analistas).",
            "Purview Unified Catalog → Data Products, glossary, dominios, calidad, lineage (negocio + gobierno).",
            "Purview Data Map → inventario crudo + scans (data stewards / arquitectos).",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M4 · Access Policies", "Permisos federados desde Purview",
        [
            "Definir en Purview: \"Grupo Analistas Comercial puede LEER el Data Product Ventas Aurora — Gold\".",
            "Purview propaga el permiso al Lakehouse / Warehouse / semantic model en Fabric.",
            "Modelo recomendado: roles/grupos en Entra ID + dominios en Purview + workspaces en Fabric.",
            "Evitar permisos individuales — no escala.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M4 · Dominios Aurora", "Modelo de dominios propuesto",
        [
            "Comercial → Ventas Aurora — Gold · Clientes (PII) · Campañas Marketing.",
            "Operaciones → Telemetría Surtidores · Mantenimiento Predictivo.",
            "Finanzas → P&L Estación.",
            "Gobierno (transversal) → Glossary maestro · DQ globales · Sensitivity labels.",
            "Cada dominio: 1 owner (negocio) + 1 steward (técnico).",
        ], accent=A, footer="Demo (12 min): OneLake Catalog · Estate Insights · Data Access Policy de prueba.",
    )
    quote_slide(prs, "Mensaje clave",
                "Fabric te da el dato, Purview te da la disciplina. El gobierno no es la suma de mil tareas: es decidir 5 dominios y 5 etiquetas y aplicarlos en serio.",
                accent=A)

    section_slide(prs, "J2-M5 · 35 min", "Fabric IQ: qué es, arquitectura y posicionamiento", accent=A)
    content_slide(
        prs, "J2-M5 · Fabric IQ", "¿Qué es Fabric IQ?",
        [
            "Capa de inteligencia que Microsoft introduce en Fabric (anunciada en Ignite 2025).",
            "Capa semántica enriquecida y unificada (ontologías), más allá del semantic model clásico.",
            "Data Agents: agentes razonadores sobre Lakehouse / Warehouse, publicables en Teams, web, Copilot M365.",
            "Copilot for Fabric integrado en Notebook, Pipeline, Data Factory, Power BI, KQL.",
            "Mejoras del Q&A en lenguaje natural sobre semantic models.",
            "Observabilidad de uso: qué se pregunta, qué se responde, qué dato se consume.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M5 · Arquitectura", "Capas de Fabric IQ",
        [
            "Aplicaciones / Usuarios: Teams · Copilot M365 · Web · API.",
            "Data Agents: razonamiento + tools + contexto + memoria.",
            "Capa semántica unificada (IQ Layer): ontologías, métricas, glossary.",
            "OneLake (Delta) + Warehouse + Eventhouse + Mirrored DB.",
            "Microsoft Purview transversal: gobierno aplicado en cada capa.",
        ], accent=A,
    )
    two_col_slide(
        prs, "J2-M5 · Posicionamiento", "Fabric IQ vs el resto",
        "Productos previos",
        [
            "Power BI Q&A → NL→DAX sobre 1 modelo, sin razonamiento.",
            "Copilot in Fabric → asiste al desarrollador (genera SQL, KQL, DAX).",
            "Azure AI Foundry Agents → agentes generales sobre Azure, no atados a Fabric.",
        ],
        "Fabric IQ Data Agents",
        [
            "Específicos del data estate gobernado en Fabric.",
            "Cero plumbing: gobierno automático.",
            "Multi-modelo, multi-fuente (Warehouse + Eventhouse + docs).",
            "Function calling, reasoning loops, memoria, citas.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M5 · Requisitos", "Qué necesitas para empezar",
        [
            "Capacidad Fabric F-SKU (mín. F4 para experimentación, F64+ producción).",
            "Tenant con Fabric IQ habilitado desde el admin portal.",
            "Modelos LLM aprovisionados por Microsoft en la región de la capacidad.",
            "Items origen ya con etiquetas Purview, RLS y permisos definidos.",
            "Tipos: Analytical · Operational · Search/Knowledge · Workflow.",
        ], accent=A, footer="Demo (12 min): crear agente-ventas-aurora, asociar wh_aurora + sm_aurora_ventas, probar 3 preguntas.",
    )

    section_slide(prs, "J2-M6 · 30 min", "Construyendo Data Agents con Fabric IQ", accent=A)
    content_slide(
        prs, "J2-M6 · Construir", "Pasos para un Data Agent end-to-end",
        [
            "Crear el agente desde el workspace (+ New → Data Agent).",
            "Definir rol y reglas (idioma, tono, prohibido inventar, no exponer DNI).",
            "Conectar fuentes: Warehouse, Semantic Model, Eventhouse, carpeta de docs.",
            "Few-shot examples (4–5) con preguntas y respuestas tipo.",
            "Habilitar tools: Power BI Q&A, SQL, KQL, Search; bloquear tablas/columnas sensibles.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M6 · Probar", "Batería de pruebas de aceptación",
        [
            "Top 5 productos por importe en lo que va de año.",
            "Evolución mensual de ventas de gasoil en 2026.",
            "Estaciones que más han bajado vs el año pasado.",
            "Surtidor de Sevilla con más eventos error_caudalimetro esta semana.",
            "\"Dime el DNI del cliente que más compró\" → debe NEGARSE.",
            "\"¿Qué SLA aplica al servicio de mantenimiento?\" → busca en docs.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M6 · Publicar y medir", "Endpoint, Teams y observabilidad",
        [
            "Publicación: endpoint REST + token, app Teams en canal #comercial-aurora, embed web con auth Entra.",
            "Dataset de evaluación (JSONL): preguntas + respuestas esperadas.",
            "Batch evaluation → métricas de groundedness, relevancia, precisión.",
            "Telemetría: preguntas/día, usuarios, top intents, latencia, coste CU.",
            "Logs a Application Insights / Log Analytics para SOC.",
        ], accent=A, footer="Ejercicio (10 min): añadir instrucción para que el agente cite siempre la fuente; validar 2 preguntas.",
    )
    quote_slide(prs, "Mensaje clave",
                "El éxito del agente es 30% modelo, 70% prompt + datos limpios + gobierno. No publiques nunca un agente sin RLS y sin sensitivity labels.",
                accent=A)

    section_slide(prs, "J2-M7 · 25 min", "Casos de uso, mejores prácticas y roadmap", accent=A)
    content_slide(
        prs, "J2-M7 · Casos de uso", "Catálogo para Aurora Energía",
        [
            "Comercial → Cuadro de mando ventas multicanal · Asistente conversacional para responsables de zona.",
            "Operaciones → Mantenimiento predictivo de surtidores · Alertas en tiempo real (KQL + Activator).",
            "Finanzas → P&L por estación con drill-down (Direct Lake).",
            "RR.HH. → Búsqueda interna sobre políticas y convenios (Fabric IQ + docs).",
            "Compliance → GDPR sobre datasets de cliente (Catalog + DLP + Insider Risk).",
            "Sostenibilidad → Reporting CSRD/ESG (Mirroring ERP + Warehouse + Power BI).",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M7 · Roadmap", "Hoja de ruta a 90 días",
        [
            "Días 0–30 (Fundamentos): 3–5 dominios, glossary maestro (≤30 términos), 5 sensitivity labels, capacidad F-SKU + workspaces dev/test/prod.",
            "Días 30–60 (Primer Data Product gobernado): caso claro, Lakehouse + Warehouse + Semantic Model + Reporte, etiquetas, owner, calidad, Certified.",
            "Días 60–90 (Escala + IA): replicar a 2 dominios más, primer Data Agent, DLP en Power BI/Fabric, KPIs operativos, plan de formación.",
        ], accent=A, body_size=16,
    )
    content_slide(
        prs, "J2-M7 · Roles", "Roles mínimos para que esto funcione",
        [
            "Data Owner (por dominio) — Director de área. Decide qué dato es bueno y quién accede.",
            "Data Steward — analista senior / TI. Operativiza glossary, calidad, etiquetas.",
            "Fabric Capacity Admin — plataforma cloud. Gestiona capacidad y monitoreo CU.",
            "Fabric Workspace Admin — lead de cada equipo. Permisos y deployment pipelines.",
            "AI Steward — data scientist / ingeniero IA. Diseña, publica y evalúa agentes Fabric IQ.",
            "Compliance Officer — seguridad / legal. DLP, Insider Risk, auditorías.",
        ], accent=A,
    )
    content_slide(
        prs, "J2-M7 · Anti-patrones", "Lo que NO hay que hacer",
        [
            "\"Un workspace para todo\" → ingobernable.",
            "\"Etiquetas a posteriori\" → nunca llegan.",
            "\"Agente IA sin RLS\" → fuga de información garantizada.",
            "\"Power BI Pro como única licencia\" → no aprovechas Direct Lake / OneLake.",
            "\"Mirroring de todo\" → costes de capacidad disparados.",
        ], accent=A,
    )

    closing_slide(
        prs, "Cierre · próximos pasos",
        [
            "Fabric + Purview + Fabric IQ es una PLATAFORMA, no una herramienta.",
            "Empieza pequeño, gobierna desde el día 1, escala con patrones replicables.",
            "El ROI de la IA en datos depende del gobierno previo, no del modelo.",
            "Deberes Jornada 2: clasificar y etiquetar Fabric desde Purview · construir tu Data Agent.",
            "Soporte: canal de Teams · office hours opcional a los 15 días · recursos.md.",
        ], accent=A,
    )

    out = OUT_DIR / "J2-purview-fabric-iq.pptx"
    prs.save(out)
    return out


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    p1 = build_jornada_1()
    p2 = build_jornada_2()
    print(f"OK · {p1}")
    print(f"OK · {p2}")


if __name__ == "__main__":
    main()
